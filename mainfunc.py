import streamlit as st
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv

load_dotenv()
os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def get_docx_text(docx_docs):
    text = ""
    for docx in docx_docs:
        doc = Document(docx)
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
    return text

def get_pptx_text(pptx_docs):
    text = ""
    for pptx in pptx_docs:
        presentation = Presentation(pptx)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text

def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    return vector_store

def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context, make sure to provide all the details, if the answer is not in
    provided context just say, "answer is not available in the context", don't provide the wrong answer\n\n
    Context:\n {context}?\n
    Question: \n{question}\n

    Answer:
    """

    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)

    return chain

def process_files(uploaded_files):
    raw_text = ""
    pdf_docs = [file for file in uploaded_files if file.name.endswith('.pdf')]
    docx_docs = [file for file in uploaded_files if file.name.endswith('.docx')]
    pptx_docs = [file for file in uploaded_files if file.name.endswith('.pptx')]

    if pdf_docs:
        raw_text += get_pdf_text(pdf_docs)
    if docx_docs:
        raw_text += get_docx_text(docx_docs)
    if pptx_docs:
        raw_text += get_pptx_text(pptx_docs)

    if raw_text:
        text_chunks = get_text_chunks(raw_text)
        vector_store = get_vector_store(text_chunks)
        st.session_state.vector_store = vector_store
        st.session_state.chain = get_conversational_chain()
        st.success("Processing completed.")
    else:
        st.error("No valid files uploaded.")

def user_input(user_question):
    vector_store = st.session_state.vector_store
    docs = vector_store.similarity_search(user_question)

    chain = st.session_state.chain
    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)

    print(response)
    st.write("Reply: ", response["output_text"])

def main():
    st.set_page_config(page_title="Chat with Documents")
    st.header("Chat with PDF, DOCX, PPTX using Gemini")

    if "vector_store" not in st.session_state:
        st.session_state.vector_store = None
        st.session_state.chain = None

    with st.sidebar:
        st.title("Menu:")
        uploaded_files = st.file_uploader("Upload your PDF, DOCX, and PPTX Files", accept_multiple_files=True, type=['pdf', 'docx', 'pptx'])
        if st.button("Submit & Process"):
            with st.spinner("Processing..."):
                process_files(uploaded_files)

    if st.session_state.vector_store and st.session_state.chain:
        user_question = st.text_input("Ask a Question from the Files")
        if user_question:
            user_input(user_question)

if __name__ == "__main__":
    main()
